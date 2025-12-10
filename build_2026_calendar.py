# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lunar_python import Solar
import calendar

YEAR = 2027

# ---------- 基本外观配置 ----------
BG_COLOR = RGBColor(250, 250, 250)
TITLE_COLOR = RGBColor(40, 40, 40)
WEEK_HEADER_COLOR = RGBColor(60, 60, 60)
DAY_COLOR = RGBColor(30, 30, 30)
WEEKEND_COLOR = RGBColor(200, 0, 0)
LABEL_COLOR = RGBColor(220, 60, 20)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ---------- 封面 ----------
cover = prs.slides.add_slide(prs.slide_layouts[6])
bg = cover.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(15, 18, 30)

tb = cover.shapes.add_textbox(Cm(3), Cm(3), Cm(24), Cm(4))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = f"{YEAR} 年 日历（公历 + 农历 + 节气 + 节日）"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 215, 0)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "自动生成 · 可编辑 · 适合 PPT 展示 / 打印"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(220, 220, 220)
p2.alignment = PP_ALIGN.CENTER

# ---------- 固定：24 节气（用 Solar 转 Lunar 去判断 term） ----------
# lunar_python 自己会算节气，这里用标签方式再补充一个“显式节气表”也行
# 但直接用 lunar.getTerm() 就能拿到当日节气名（无则 None）

# 一些常见节日（按农历 / 公历混合）
def get_festival_label(solar):
    """返回节日名称列表（用 lunar_python 自带的节日 + 少量自定义补充）"""
    labels = []

    # 1）可选的公历固定节日（补充，避免有些版本没内置）
    y, m, d = solar.getYear(), solar.getMonth(), solar.getDay()
    if m == 1 and d == 1:
        labels.append("元旦")
    if m == 5 and d == 1:
        labels.append("劳动节")
    if m == 6 and d == 1:
        labels.append("儿童节")
    if m == 9 and d == 10:
        labels.append("教师节")
    if m == 10 and d == 1:
        labels.append("国庆节")

    # 2）使用 lunar_python 自带的节日（包含农历节日、部分公历节日）
    try:
        fest_list = solar.getFestivals()  # 返回 list[str]，如 ["春节"]、["中秋节"] 等
        for f in fest_list:
            labels.append(f)
    except AttributeError:
        # 某些极老版本如果没有 getFestivals，可以忽略这一步
        pass

    # 去重，保留顺序
    unique = []
    for x in labels:
        if x not in unique:
            unique.append(x)
    return unique

    """返回节日名称（可按需扩展）"""
    y, m, d = solar.getYear(), solar.getMonth(), solar.getDay()
    lunar = solar.getLunar()
    lm, ld = lunar.getMonth(), lunar.getDay()

    labels = []

    # 公历节日
    if m == 1 and d == 1:
        labels.append("元旦")
    if m == 5 and d == 1:
        labels.append("劳动节")
    if m == 6 and d == 1:
        labels.append("儿童节")
    if m == 9 and d == 10:
        labels.append("教师节")
    if m == 10 and d == 1:
        labels.append("国庆节")

    # 农历节日
    # 春节：正月初一
    if lm == 1 and ld == 1:
        labels.append("春节")
    # 元宵：正月十五
    if lm == 1 and ld == 15:
        labels.append("元宵节")
    # 端午：五月初五
    if lm == 5 and ld == 5:
        labels.append("端午节")
    # 七夕：七月初七
    if lm == 7 and ld == 7:
        labels.append("七夕节")
    # 中元：七月十五
    if lm == 7 and ld == 15:
        labels.append("中元节")
    # 中秋：八月十五
    if lm == 8 and ld == 15:
        labels.append("中秋节")
    # 重阳：九月初九
    if lm == 9 and ld == 9:
        labels.append("重阳节")
    # 腊八：腊月初八
    if lunar.isLunarMonthLeap():
        pass  # 不考虑闰月腊八
    if lm == 12 and ld == 8:
        labels.append("腊八节")
    # 除夕：该农历年的最后一天
    if lunar.isLastDay():
        labels.append("除夕")

    return labels

week_headers = ["一", "二", "三", "四", "五", "六", "日"]

for month in range(1, 13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_COLOR

    # 标题
    tb = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(26), Cm(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{YEAR} 年 {month} 月"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    # 日历表格
    rows, cols = 7, 7  # 1 行星期 + 6 行日期
    left, top, width, height = Cm(1.5), Cm(3), Cm(26), Cm(13)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(cols):
        table.columns[i].width = width // cols

    # 星期标题行
    for i, wd in enumerate(week_headers):
        cell = table.cell(0, i)
        cell.text = wd
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = WEEK_HEADER_COLOR

    cal = calendar.Calendar(firstweekday=0)  # 周一=0

    for week_idx, week in enumerate(cal.monthdayscalendar(YEAR, month)):
        for weekday_idx, day_num in enumerate(week):
            if day_num == 0:
                continue
            row = week_idx + 1
            col = weekday_idx
            cell = table.cell(row, col)
            tf = cell.text_frame
            tf.clear()

            # 公历
            solar = Solar.fromYmd(YEAR, month, day_num)
            lunar = solar.getLunar()

            # 1）公历日期（大字）
            p_day = tf.add_paragraph()
            p_day.alignment = PP_ALIGN.LEFT
            r_day = p_day.add_run()
            r_day.text = str(day_num)
            r_day.font.size = Pt(20)
            r_day.font.bold = True
            # 周末变红
            if weekday_idx >= 5:
                r_day.font.color.rgb = WEEKEND_COLOR
            else:
                r_day.font.color.rgb = DAY_COLOR

            # 2）农历日名（小字）：如 初一、初二、廿三
            p_lunar = tf.add_paragraph()
            p_lunar.alignment = PP_ALIGN.LEFT
            r_lunar = p_lunar.add_run()
            r_lunar.text = lunar.getDayInChinese()  # 只要日名，不带“正月”等
            r_lunar.font.size = Pt(11)
            r_lunar.font.color.rgb = RGBColor(120, 120, 120)

            # 3）节气 + 节日
            labels = []

            # lunar_python 的接口是 getJieQi()，不是 getTerm()
            term = lunar.getJieQi()  # 若当天是节气，返回节气名（例如“雨水”），否则返回空字符串
            if term:
                labels.append(term)

            labels.extend(get_festival_label(solar))

            if labels:
                p_lab = tf.add_paragraph()
                p_lab.alignment = PP_ALIGN.LEFT
                r_lab = p_lab.add_run()
                r_lab.text = " / ".join(labels)
                r_lab.font.size = Pt(11)
                r_lab.font.color.rgb = LABEL_COLOR

# 保存
out_file = f"{YEAR}_全年日历_公历_农历_节气_节日版.pptx"
prs.save(out_file)
print("已生成：", out_file)
